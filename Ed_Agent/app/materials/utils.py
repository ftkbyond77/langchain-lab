import os
from pdfminer.high_level import extract_text
from pptx import Presentation

def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    if ext == ".pdf":
        text = extract_text(file_path)
    elif ext in [".ppt", ".pptx"]:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    else:
        # Plain text
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read()
    return text
